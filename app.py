"""
Editorial QA Agent — Gradio Web UI
====================================
Customer-facing front-end for the PPTX Editorial QA pipeline.

Features:
  • One-click folder selection via native Windows dialog
  • Automatically reviews ALL PPTX decks in the selected folder
  • Per-deck and per-slide findings with high-contrast design
  • Summary dashboard with severity/category breakdowns
  • Full agent report

Launch:  python app.py
"""

import os
import re
import json
import sys
import subprocess
import time
import logging
import pathlib
from dataclasses import dataclass
from typing import List, Dict, Any

import gradio as gr
from dotenv import load_dotenv
from pptx import Presentation as PptxPresentation

# Azure AI Foundry SDK v2
from azure.identity import AzureCliCredential
from azure.ai.projects import AIProjectClient
from azure.ai.projects.models import PromptAgentDefinition, FunctionTool
from openai.types.responses.response_input_param import FunctionCallOutput

# ═══════════════════════════════════════════════════════════════════════════════
# Configuration
# ═══════════════════════════════════════════════════════════════════════════════
PROJECT_ROOT = pathlib.Path(__file__).parent
load_dotenv(str(PROJECT_ROOT / ".env"), override=True)

ENDPOINT = os.environ["AZURE_AI_PROJECT_ENDPOINT"]
DEFAULT_MODEL = os.environ.get("AZURE_AI_MODEL_DEPLOYMENT_NAME", "gpt-4o")

logging.basicConfig(level=logging.INFO, format="%(asctime)s [%(levelname)s] %(message)s")
logger = logging.getLogger("editorial_qa_ui")


# ═══════════════════════════════════════════════════════════════════════════════
# Core extraction / chunking
# ═══════════════════════════════════════════════════════════════════════════════

@dataclass
class SlideBlock:
    slide_num: int
    markdown: str


def _clean(s: str) -> str:
    s = s.replace("\xa0", " ")
    s = re.sub(r"[ ]+", " ", s)
    s = re.sub(r"\n{3,}", "\n\n", s)
    return s.strip()


def pptx_to_markdown_slides(pptx_path: str) -> List[SlideBlock]:
    prs = PptxPresentation(pptx_path)
    blocks: List[SlideBlock] = []
    for i, slide in enumerate(prs.slides, start=1):
        lines = [f"# Slide {i}"]
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue
            tf = shape.text_frame
            if not tf or not tf.text:
                continue
            for paragraph in tf.paragraphs:
                txt = _clean(paragraph.text)
                if not txt:
                    continue
                if paragraph.level and paragraph.level > 0:
                    indent = "  " * (paragraph.level - 1)
                    lines.append(f"{indent}- {txt}")
                else:
                    lines.append(txt)
        md = _clean("\n".join(lines))
        if md and md != f"# Slide {i}":
            blocks.append(SlideBlock(i, md))
    return blocks


def join_slides(blocks: List[SlideBlock]) -> str:
    return "\n\n\n".join(b.markdown for b in blocks)


# ═══════════════════════════════════════════════════════════════════════════════
# Visual / layout metadata extraction
# ═══════════════════════════════════════════════════════════════════════════════

def _emu_to_inches(emu) -> float:
    """Convert EMU (English Metric Units) to inches. 914400 EMU = 1 inch."""
    if emu is None:
        return 0.0
    return round(emu / 914400, 2)


def _color_hex(color_obj) -> str | None:
    """Safely extract hex color string from a pptx color object."""
    try:
        if color_obj and color_obj.type is not None:
            return str(color_obj.rgb) if hasattr(color_obj, 'rgb') and color_obj.rgb else None
    except Exception:
        pass
    return None


def _font_props(font) -> dict:
    """Extract font properties from a pptx font object."""
    props = {}
    try:
        if font.name:
            props["name"] = font.name
        if font.size:
            props["size_pt"] = round(font.size.pt, 1)
        if font.bold is not None:
            props["bold"] = font.bold
        if font.italic is not None:
            props["italic"] = font.italic
        c = _color_hex(font.color)
        if c:
            props["color"] = c
    except Exception:
        pass
    return props


def extract_slide_visual_metadata(pptx_path: str) -> List[dict]:
    """
    Extract visual/layout metadata for every slide in a PPTX file.

    For each shape on each slide, captures:
      - position (left, top in inches), size (width, height in inches)
      - shape type, rotation
      - font properties per text run (name, size, bold, italic, color)
      - fill color (solid fills only)

    Returns a list of per-slide metadata dicts.
    """
    prs = PptxPresentation(pptx_path)
    slides_meta = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        shapes_meta = []
        for shape in slide.shapes:
            shape_info = {
                "shape_type": str(shape.shape_type) if shape.shape_type else "unknown",
                "name": shape.name or "",
                "left_in": _emu_to_inches(shape.left),
                "top_in": _emu_to_inches(shape.top),
                "width_in": _emu_to_inches(shape.width),
                "height_in": _emu_to_inches(shape.height),
                "rotation": shape.rotation if shape.rotation else 0,
            }

            # Fill color
            try:
                fill = shape.fill
                if fill and fill.type is not None:
                    fc = _color_hex(fill.fore_color)
                    if fc:
                        shape_info["fill_color"] = fc
            except Exception:
                pass

            # Text / font properties
            if hasattr(shape, "text_frame") and shape.text_frame:
                paras = []
                for para in shape.text_frame.paragraphs:
                    txt = _clean(para.text)
                    if not txt:
                        continue
                    # Collect unique font specs across runs
                    fonts_in_para = []
                    for run in para.runs:
                        fp = _font_props(run.font)
                        if fp and fp not in fonts_in_para:
                            fonts_in_para.append(fp)
                    para_info = {"text": txt[:120]}
                    if fonts_in_para:
                        para_info["fonts"] = fonts_in_para
                    paras.append(para_info)
                if paras:
                    shape_info["paragraphs"] = paras

            shapes_meta.append(shape_info)

        slides_meta.append({
            "slide_num": slide_idx,
            "shape_count": len(shapes_meta),
            "shapes": shapes_meta,
        })

    return slides_meta


SINGLE_PASS_TOKEN_THRESHOLD = 90_000
DEFAULT_WINDOW_SIZE = 1200
DEFAULT_OVERLAP = 250


def estimate_token_count(text: str) -> int:
    return int(len(text.split()) * 1.3)


def sliding_window_chunks(text, window_size=DEFAULT_WINDOW_SIZE, overlap=DEFAULT_OVERLAP):
    words = text.split()
    chunks, start = [], 0
    while start < len(words):
        end = min(start + window_size, len(words))
        chunks.append((start, end, " ".join(words[start:end])))
        if end == len(words):
            break
        start = max(0, end - overlap)
    return chunks


def adaptive_chunk(deck_md, window_size=DEFAULT_WINDOW_SIZE, overlap=DEFAULT_OVERLAP):
    tokens = estimate_token_count(deck_md)
    words = deck_md.split()
    if tokens < SINGLE_PASS_TOKEN_THRESHOLD:
        return [(0, len(words), deck_md)]
    return sliding_window_chunks(deck_md, window_size, overlap)


# ═══════════════════════════════════════════════════════════════════════════════
# Function tools (shared state per session via closures)
# ═══════════════════════════════════════════════════════════════════════════════
findings_accumulator: List[Dict[str, Any]] = []
extracted_deck_cache: Dict[str, str] = {}
SEVERITY_RANK = {"Critical": 4, "High": 3, "Medium": 2, "Low": 1}


def extract_deck(pptx_path: str) -> str:
    blocks = pptx_to_markdown_slides(pptx_path)
    deck_md = join_slides(blocks)
    extracted_deck_cache[pptx_path] = deck_md
    return json.dumps({
        "slide_count": len(blocks),
        "word_count": len(deck_md.split()),
        "estimated_tokens": estimate_token_count(deck_md),
        "markdown": deck_md,
    })


def get_review_windows(markdown: str, window_size: int = 1200, overlap: int = 250) -> str:
    chunks = adaptive_chunk(markdown, window_size, overlap)
    result = []
    for idx, (start, end, text) in enumerate(chunks):
        slide_nums = [int(m) for m in re.findall(r"# Slide (\d+)", text)]
        result.append({
            "chunk_id": idx, "start_word": start, "end_word": end,
            "word_count": end - start, "slides_covered": slide_nums, "text": text,
        })
    return json.dumps({"total_chunks": len(result),
                        "mode": "single_pass" if len(result) == 1 else "windowed",
                        "chunks": result})


def store_chunk_findings(chunk_id: int, findings_json: str) -> str:
    try:
        findings = json.loads(findings_json)
        if not isinstance(findings, list):
            findings = [findings]
    except json.JSONDecodeError:
        return json.dumps({"error": "Invalid JSON in findings", "stored": False})
    for f in findings:
        f["source_chunk"] = chunk_id
    logger.info(f"store_chunk_findings(chunk_id={chunk_id}): {len(findings)} findings")
    for f in findings:
        logger.info(f"  -> {f.get('slides','?')} | {f.get('category','?')} | {f.get('issue','?')[:80]}")
    findings_accumulator.extend(findings)
    return json.dumps({"stored": True, "chunk_id": chunk_id,
                        "findings_in_chunk": len(findings),
                        "total_findings_so_far": len(findings_accumulator)})


def merge_and_dedupe_findings() -> str:
    if not findings_accumulator:
        return json.dumps({"merged_findings": [], "total": 0, "duplicates_removed": 0})
    seen: Dict[str, Dict[str, Any]] = {}
    for finding in findings_accumulator:
        slides_key = str(finding.get("slides", finding.get("slide", ""))).strip().lower()
        issue_key = str(finding.get("issue", finding.get("category", ""))).strip().lower()
        evidence_key = str(finding.get("evidence", finding.get("remediation", ""))).strip().lower()
        dedup_key = f"{slides_key}|{issue_key}|{evidence_key}"
        current_rank = SEVERITY_RANK.get(str(finding.get("flag", finding.get("severity", "Low"))), 0)
        if dedup_key in seen:
            existing_rank = SEVERITY_RANK.get(
                str(seen[dedup_key].get("flag", seen[dedup_key].get("severity", "Low"))), 0)
            if current_rank > existing_rank:
                seen[dedup_key] = finding
        else:
            seen[dedup_key] = finding
    merged = list(seen.values())

    def sort_key(f):
        sev = SEVERITY_RANK.get(str(f.get("flag", f.get("severity", "Low"))), 0)
        slide_str = str(f.get("slides", f.get("slide", "999")))
        slide_num = int(re.search(r"\d+", slide_str).group()) if re.search(r"\d+", slide_str) else 999
        return (-sev, slide_num)

    merged.sort(key=sort_key)
    return json.dumps({"merged_findings": merged, "total": len(merged),
                        "duplicates_removed": len(findings_accumulator) - len(merged)})


def extract_deck_visual(pptx_path: str) -> str:
    """Extract visual/layout metadata (fonts, colors, positions, sizes) for all slides."""
    try:
        meta = extract_slide_visual_metadata(pptx_path)
        return json.dumps({"slide_count": len(meta), "slides": meta})
    except Exception as e:
        logger.error(f"Visual extraction failed: {e}")
        return json.dumps({"error": str(e)})


def analyze_visual_consistency(pptx_path: str) -> str:
    """
    Deterministic visual consistency analysis.

    Extracts visual metadata, computes the dominant pattern for titles and body,
    then returns ONLY the deviations — pre-computed so the LLM doesn't have to
    scan thousands of tokens of raw shape data.
    """
    from collections import Counter
    meta = extract_slide_visual_metadata(pptx_path)

    # ── Classify shapes per slide: title = first text shape, body = subsequent ──
    slides_info = []
    for s in meta:
        text_shapes = [sh for sh in s["shapes"] if "paragraphs" in sh]
        deco_shapes = [sh for sh in s["shapes"] if "paragraphs" not in sh and sh.get("fill_color")]
        title_sh = text_shapes[0] if text_shapes else None
        body_shs = text_shapes[1:] if len(text_shapes) > 1 else []
        slides_info.append({
            "slide_num": s["slide_num"],
            "shape_count": s["shape_count"],
            "title": title_sh,
            "body": body_shs,
            "deco": deco_shapes,
            "all_shapes": s["shapes"],
        })

    def _first_font(shape):
        if not shape or "paragraphs" not in shape:
            return {}
        for p in shape["paragraphs"]:
            if "fonts" in p and p["fonts"]:
                return p["fonts"][0]
        return {}

    # ── Compute dominant title pattern ──
    title_fonts = Counter()
    title_sizes = Counter()
    title_bolds = Counter()
    title_italics = Counter()
    title_colors = Counter()
    title_lefts = Counter()
    title_tops = Counter()
    title_widths = Counter()

    for si in slides_info:
        f = _first_font(si["title"])
        if not f:
            continue
        if f.get("name"):   title_fonts[f["name"]] += 1
        if f.get("size_pt"): title_sizes[f["size_pt"]] += 1
        if "bold" in f:      title_bolds[f["bold"]] += 1
        if "italic" in f:    title_italics[f["italic"]] += 1
        if f.get("color"):   title_colors[f["color"]] += 1
        if si["title"]:
            title_lefts[si["title"].get("left_in")] += 1
            title_tops[si["title"].get("top_in")] += 1
            title_widths[si["title"].get("width_in")] += 1

    dom_t = {
        "font": title_fonts.most_common(1)[0][0] if title_fonts else None,
        "size": title_sizes.most_common(1)[0][0] if title_sizes else None,
        "bold": title_bolds.most_common(1)[0][0] if title_bolds else None,
        "italic": title_italics.most_common(1)[0][0] if title_italics else None,
        "color": title_colors.most_common(1)[0][0] if title_colors else None,
        "left": title_lefts.most_common(1)[0][0] if title_lefts else None,
        "top": title_tops.most_common(1)[0][0] if title_tops else None,
        "width": title_widths.most_common(1)[0][0] if title_widths else None,
    }

    # ── Compute dominant body pattern ──
    body_fonts = Counter()
    body_sizes = Counter()
    body_bolds = Counter()
    body_italics = Counter()
    body_colors = Counter()

    for si in slides_info:
        for bsh in si["body"]:
            f = _first_font(bsh)
            if not f:
                continue
            if f.get("name"):   body_fonts[f["name"]] += 1
            if f.get("size_pt"): body_sizes[f["size_pt"]] += 1
            if "bold" in f:      body_bolds[f["bold"]] += 1
            if "italic" in f:    body_italics[f["italic"]] += 1
            if f.get("color"):   body_colors[f["color"]] += 1

    dom_b = {
        "font": body_fonts.most_common(1)[0][0] if body_fonts else None,
        "size": body_sizes.most_common(1)[0][0] if body_sizes else None,
        "bold": body_bolds.most_common(1)[0][0] if body_bolds else None,
        "italic": body_italics.most_common(1)[0][0] if body_italics else None,
        "color": body_colors.most_common(1)[0][0] if body_colors else None,
    }

    # ── Dominant accent fill ──
    fill_colors = Counter()
    for si in slides_info:
        for d in si["deco"]:
            if d.get("fill_color"):
                fill_colors[d["fill_color"]] += 1
    dom_fill = fill_colors.most_common(1)[0][0] if fill_colors else None

    # ── Average shape count (for clutter detection) ──
    counts = [si["shape_count"] for si in slides_info]
    avg_shapes = sum(counts) / len(counts) if counts else 0

    # ══════════════════════════════════════════════════════════════════════
    # Find deviations
    # ══════════════════════════════════════════════════════════════════════
    issues = []

    for si in slides_info:
        sn = si["slide_num"]
        tf = _first_font(si["title"])
        tsh = si["title"]

        # ── Cross-slide title font deviations ──
        if tf.get("name") and dom_t["font"] and tf["name"] != dom_t["font"]:
            issues.append({"slide": sn, "type": "Visual", "element": "title",
                           "check": "font_family",
                           "found": tf["name"], "expected": dom_t["font"],
                           "detail": f"Title font '{tf['name']}' vs dominant '{dom_t['font']}'"})
        if tf.get("size_pt") and dom_t["size"] and tf["size_pt"] != dom_t["size"]:
            issues.append({"slide": sn, "type": "Visual", "element": "title",
                           "check": "font_size",
                           "found": tf["size_pt"], "expected": dom_t["size"],
                           "detail": f"Title size {tf['size_pt']}pt vs dominant {dom_t['size']}pt"})
        if "bold" in tf and dom_t["bold"] is not None and tf["bold"] != dom_t["bold"]:
            issues.append({"slide": sn, "type": "Visual", "element": "title",
                           "check": "bold",
                           "found": tf.get("bold"), "expected": dom_t["bold"],
                           "detail": f"Title bold={tf['bold']} vs dominant bold={dom_t['bold']}"})
        if "italic" in tf and dom_t["italic"] is not None and tf["italic"] != dom_t["italic"]:
            issues.append({"slide": sn, "type": "Visual", "element": "title",
                           "check": "italic",
                           "found": tf.get("italic"), "expected": dom_t["italic"],
                           "detail": f"Title italic={tf['italic']} vs dominant italic={dom_t['italic']}"})
        if tf.get("color") and dom_t["color"] and tf["color"] != dom_t["color"]:
            issues.append({"slide": sn, "type": "Visual", "element": "title",
                           "check": "text_color",
                           "found": tf["color"], "expected": dom_t["color"],
                           "detail": f"Title color #{tf['color']} vs dominant #{dom_t['color']}"})

        # ── Cross-slide title position/size deviations ──
        if tsh and dom_t["top"] is not None:
            if abs((tsh.get("top_in") or 0) - dom_t["top"]) > 0.3:
                issues.append({"slide": sn, "type": "Layout", "element": "title",
                               "check": "position_top",
                               "found": tsh.get("top_in"), "expected": dom_t["top"],
                               "detail": f"Title top={tsh.get('top_in')}\" vs dominant {dom_t['top']}\""})
        if tsh and dom_t["left"] is not None:
            if abs((tsh.get("left_in") or 0) - dom_t["left"]) > 0.3:
                issues.append({"slide": sn, "type": "Layout", "element": "title",
                               "check": "position_left",
                               "found": tsh.get("left_in"), "expected": dom_t["left"],
                               "detail": f"Title left={tsh.get('left_in')}\" vs dominant {dom_t['left']}\""})
        if tsh and dom_t["width"] is not None:
            if abs((tsh.get("width_in") or 0) - dom_t["width"]) > 1.0:
                issues.append({"slide": sn, "type": "Visual", "element": "title",
                               "check": "shape_width",
                               "found": tsh.get("width_in"), "expected": dom_t["width"],
                               "detail": f"Title width={tsh.get('width_in')}\" vs dominant {dom_t['width']}\""})

        # ── Cross-slide body font deviations (consolidated per slide) ──
        body_dev_fonts = set()
        body_dev_sizes = set()
        body_dev_bolds = set()
        body_dev_italics = set()
        body_dev_colors = set()
        for bsh in si["body"]:
            bf = _first_font(bsh)
            if bf.get("name") and dom_b["font"] and bf["name"] != dom_b["font"]:
                body_dev_fonts.add(bf["name"])
            if bf.get("size_pt") and dom_b["size"] and bf["size_pt"] != dom_b["size"]:
                body_dev_sizes.add(bf["size_pt"])
            if "bold" in bf and dom_b["bold"] is not None and bf["bold"] != dom_b["bold"]:
                body_dev_bolds.add(bf["bold"])
            if "italic" in bf and dom_b["italic"] is not None and bf["italic"] != dom_b["italic"]:
                body_dev_italics.add(bf["italic"])
            if bf.get("color") and dom_b["color"] and bf["color"] != dom_b["color"]:
                body_dev_colors.add(bf["color"])
        if body_dev_fonts:
            issues.append({"slide": sn, "type": "Visual", "element": "body",
                           "check": "font_family",
                           "found": sorted(body_dev_fonts), "expected": dom_b["font"],
                           "detail": f"Body font(s) {sorted(body_dev_fonts)} vs dominant '{dom_b['font']}'"})
        if body_dev_sizes:
            issues.append({"slide": sn, "type": "Visual", "element": "body",
                           "check": "font_size",
                           "found": sorted(body_dev_sizes), "expected": dom_b["size"],
                           "detail": f"Body size(s) {sorted(body_dev_sizes)}pt vs dominant {dom_b['size']}pt"})
        if body_dev_bolds:
            issues.append({"slide": sn, "type": "Visual", "element": "body",
                           "check": "bold",
                           "found": sorted(body_dev_bolds), "expected": dom_b["bold"],
                           "detail": f"Body bold deviation vs dominant bold={dom_b['bold']}"})
        if body_dev_italics:
            issues.append({"slide": sn, "type": "Visual", "element": "body",
                           "check": "italic",
                           "found": sorted(body_dev_italics), "expected": dom_b["italic"],
                           "detail": f"Body italic deviation vs dominant italic={dom_b['italic']}"})
        if body_dev_colors:
            issues.append({"slide": sn, "type": "Visual", "element": "body",
                           "check": "text_color",
                           "found": sorted(body_dev_colors), "expected": dom_b["color"],
                           "detail": f"Body color(s) {['#'+c for c in sorted(body_dev_colors)]} vs dominant #{dom_b['color']}"})

        # ── Within-slide font/size/color mixing among body shapes ──
        if len(si["body"]) >= 2:
            fonts_on_slide = set()
            sizes_on_slide = set()
            colors_on_slide = set()
            for bsh in si["body"]:
                bf = _first_font(bsh)
                if bf.get("name"):   fonts_on_slide.add(bf["name"])
                if bf.get("size_pt"): sizes_on_slide.add(bf["size_pt"])
                if bf.get("color"):  colors_on_slide.add(bf["color"])
            if len(fonts_on_slide) > 1:
                issues.append({"slide": sn, "type": "Visual", "element": "body",
                               "check": "within_slide_font_mix",
                               "found": sorted(fonts_on_slide), "expected": "same font",
                               "detail": f"Within-slide font mixing: {sorted(fonts_on_slide)}"})
            if len(sizes_on_slide) > 1:
                issues.append({"slide": sn, "type": "Visual", "element": "body",
                               "check": "within_slide_size_mix",
                               "found": sorted(sizes_on_slide), "expected": "same size",
                               "detail": f"Within-slide size mixing: {sorted(sizes_on_slide)}pt"})
            if len(colors_on_slide) > 1:
                issues.append({"slide": sn, "type": "Visual", "element": "body",
                               "check": "within_slide_color_mix",
                               "found": sorted(colors_on_slide), "expected": "same color",
                               "detail": f"Within-slide color mixing: {['#'+c for c in sorted(colors_on_slide)]}"})

        # ── Accent fill color deviation ──
        for d in si["deco"]:
            if d.get("fill_color") and dom_fill and d["fill_color"] != dom_fill:
                issues.append({"slide": sn, "type": "Visual", "element": "accent",
                               "check": "fill_color",
                               "found": d["fill_color"], "expected": dom_fill,
                               "detail": f"Accent fill #{d['fill_color']} vs dominant #{dom_fill}"})

        # ── Layout: Clutter ──
        if si["shape_count"] > avg_shapes * 2 and si["shape_count"] > 6:
            issues.append({"slide": sn, "type": "Layout", "element": "slide",
                           "check": "clutter",
                           "found": si["shape_count"], "expected": f"avg {avg_shapes:.0f}",
                           "detail": f"{si['shape_count']} shapes vs deck avg {avg_shapes:.0f}"})

        # ── Layout: Overlap detection ──
        all_shapes = ([si["title"]] if si["title"] else []) + si["body"]
        for i_a in range(len(all_shapes)):
            for i_b in range(i_a + 1, len(all_shapes)):
                a, b = all_shapes[i_a], all_shapes[i_b]
                a_r = a.get("left_in", 0) + a.get("width_in", 0)
                a_bot = a.get("top_in", 0) + a.get("height_in", 0)
                b_r = b.get("left_in", 0) + b.get("width_in", 0)
                b_bot = b.get("top_in", 0) + b.get("height_in", 0)
                if (a.get("left_in", 0) < b_r and b.get("left_in", 0) < a_r and
                    a.get("top_in", 0) < b_bot and b.get("top_in", 0) < a_bot):
                    issues.append({"slide": sn, "type": "Layout", "element": "shapes",
                                   "check": "overlap",
                                   "found": f"{a.get('name','')} ↔ {b.get('name','')}",
                                   "expected": "no overlap",
                                   "detail": f"Shapes '{a.get('name','')}' and '{b.get('name','')}' overlap"})

        # ── Layout: Crammed margins ──
        check_shapes = ([si["title"]] if si["title"] else []) + si["body"]
        for sh in check_shapes:
            if not sh:
                continue
            l = sh.get("left_in", 1)
            t = sh.get("top_in", 1)
            if l < 0.3 or t < 0.3:
                issues.append({"slide": sn, "type": "Layout", "element": sh.get("name", "shape"),
                               "check": "crammed_margin",
                               "found": f"left={l}\", top={t}\"", "expected": ">0.3\" margin",
                               "detail": f"Shape '{sh.get('name','')}' at ({l}\", {t}\") — crammed to edge"})

        # ── Layout: Misaligned columns ──
        if len(si["body"]) >= 3:
            tops = [bsh.get("top_in", 0) for bsh in si["body"]]
            if max(tops) - min(tops) > 0.3:
                issues.append({"slide": sn, "type": "Layout", "element": "columns",
                               "check": "misaligned_columns",
                               "found": tops, "expected": "same top values",
                               "detail": f"Column tops vary: {tops} (diff={max(tops)-min(tops):.1f}\")"})

        # ── Layout: Rotation ──
        for sh in si["all_shapes"]:
            if sh.get("rotation") and sh["rotation"] != 0:
                issues.append({"slide": sn, "type": "Layout", "element": sh.get("name", "shape"),
                               "check": "rotation",
                               "found": sh["rotation"], "expected": 0,
                               "detail": f"Shape '{sh.get('name','')}' rotated {sh['rotation']}°"})

    # ══════════════════════════════════════════════════════════════════════
    # Filter: skip slide 1 (cover) and last slide (closing)
    # ══════════════════════════════════════════════════════════════════════
    last_slide = max(si["slide_num"] for si in slides_info) if slides_info else 0
    filtered = [i for i in issues
                if i["slide"] != 1 and i["slide"] != last_slide]

    # ── Severity assignment ──
    def _severity(iss):
        check = iss["check"]
        element = iss["element"]
        typ = iss["type"]
        if check == "overlap":
            return "Critical"
        if check in ("font_family", "font_size", "text_color") and element == "title":
            return "High"
        if check in ("bold", "italic") and element == "title":
            return "High"
        if check in ("font_family", "font_size", "text_color", "bold", "italic") and element == "body":
            return "Medium"
        if check.startswith("within_slide"):
            return "Medium"
        if check == "fill_color":
            return "Medium"
        if check in ("shape_width",):
            return "Medium"
        if check in ("clutter", "crammed_margin", "misaligned_columns"):
            return "Medium"
        if check in ("position_top", "position_left"):
            return "Medium"
        if check == "rotation":
            return "Low"
        return "Medium"

    # ── Build pre-formatted findings ready for store_chunk_findings ──
    formatted = []
    for iss in filtered:
        formatted.append({
            "slides": f"Slide {iss['slide']}",
            "evidence": iss["detail"],
            "issue": f"{iss['element'].title()} {iss['check'].replace('_', ' ')} — deviates from dominant pattern",
            "flag": _severity(iss),
            "remediation": f"Change to match dominant pattern ({iss['expected']})",
            "category": iss["type"],
        })

    return json.dumps({
        "dominant_title": dom_t,
        "dominant_body": dom_b,
        "dominant_accent_fill": dom_fill,
        "avg_shape_count": round(avg_shapes, 1),
        "total_inconsistencies": len(filtered),
        "findings_json": json.dumps(formatted),
    })


TOOL_FUNCTIONS = {
    "extract_deck": extract_deck,
    "extract_deck_visual": extract_deck_visual,
    "analyze_visual_consistency": analyze_visual_consistency,
    "get_review_windows": get_review_windows,
    "store_chunk_findings": store_chunk_findings,
    "merge_and_dedupe_findings": merge_and_dedupe_findings,
}


# ═══════════════════════════════════════════════════════════════════════════════
# Foundry Agent setup
# ═══════════════════════════════════════════════════════════════════════════════

SYSTEM_INSTRUCTIONS = """\
You are an editorial QA reviewer for high-stakes, client-facing deliverables (often PowerPoint).
Your job is to protect credibility by flagging errors and inconsistencies that a customer would notice.

You have access to function tools. Use them in this order:
1. Call extract_deck with the provided PPTX path to get the markdown text.
2. Call analyze_visual_consistency with the same PPTX path. This returns pre-computed visual/layout inconsistencies — dominant patterns and every deviation.
3. Call get_review_windows with the markdown to split it into review chunks.
4. For EACH chunk, perform the editorial review described below, then call store_chunk_findings with the results.
5. After ALL text chunks are reviewed, perform the VISUAL/LAYOUT REVIEW using the pre-computed inconsistencies from step 2, and call store_chunk_findings with chunk_id=-1 for visual findings.
6. After ALL reviews are complete, call merge_and_dedupe_findings to get the final output.

EDITORIAL REVIEW PROCESS (apply to each chunk):

Analyze only customer-visible content. Do not use speaker notes, comments, revision history, timestamps, or file metadata.

Must-flag categories (do not skip any):
1) Spelling: typos, misspellings, incorrect word forms, obviously wrong proper nouns in context.
2) Grammar: agreement, tense consistency, broken sentence structure, incorrect articles, awkward fragments.
3) Punctuation: missing punctuation, inconsistent style, broken quotes, inconsistent hyphenation (e.g. "real time" vs "real-time").
4) Terminology consistency: same concept named multiple ways across slides (capitalization, hyphenation, acronyms, product names, metrics).
5) Tone and voice consistency: shifts suggesting multiple authors or stitched content (formal vs casual, marketing hype vs neutral, inconsistent we/you/they, inconsistent certainty).
6) Visual consistency (from visual metadata): inconsistent fonts, font sizes, colors, shape positioning, or shape sizing within and across slides.
7) Layout quality (from visual metadata): cluttered slides, overlapping shapes, poor alignment, mismatched element spacing, or confusing visual hierarchy.

Constraints:
- ZERO TOLERANCE FOR FALSE POSITIVES. It is far worse to flag a non-issue than to miss a real one. When in doubt, do NOT flag it.
- Do not invent issues. Only flag what is unambiguously wrong in the extracted text.
- Every finding MUST include the exact verbatim text that is wrong in the "evidence" field (copy-paste from the markdown). If you cannot quote the exact erroneous text, do not report the finding.
- Do not flag layout, formatting, or visual design choices — you only see extracted text, not how it is rendered on the slide.
- Do not flag statistics, percentages, or numeric data as errors unless there is a clear internal contradiction (e.g., "ROI is 340%" on one slide and "ROI is 280%" on another).
- Do not flag sentence fragments or bullet-point style phrasing — PowerPoint slides are not prose documents. Bullet points, headlines, and short phrases are normal and expected. Specifically:
  * Do NOT flag bullet points or headings for missing periods, missing articles ("a", "the"), or missing conjunctions.
  * Do NOT flag short noun phrases like "Observability and agent controls" or "Enterprise-grade security" — these are standard slide labels, not broken sentences.
  * Do NOT suggest adding periods to the end of bullet points or headings.
  * Do NOT flag lack of full sentence structure in bullet items — they are meant to be concise.
- Do not flag punctuation in headings, labels, or statistics (e.g., "85%" is correct, not a punctuation error).
- Do not rewrite the document. Provide targeted remediation only.
- If a slide has no issues, do not list it.
- Only flag contradictions if both statements are visible and you can cite the slide numbers involved.

For each chunk, perform this 4-pass process:
Pass 1 (Inventory): Build a Canonical Glossary from recurring key terms, acronyms, product names, and metrics. Choose canonical form based on most common or most formal usage. Identify dominant tone in 1-2 sentences.
Pass 2 (Per-slide QA): Scan each slide for spelling, grammar, and punctuation issues. Flag tone anomalies only when clearly divergent from surrounding slides.
Pass 3 (Cross-slide consistency): Using your Canonical Glossary and dominant voice, flag terminology drift, inconsistent acronym usage, and tone drift.
Pass 4 (VERIFICATION — MANDATORY): Before storing findings, re-read each finding and ask: "Can I point to the EXACT wrong text and explain precisely why it is wrong?" Drop any finding where the answer is no. Drop any finding that is subjective, stylistic preference, or layout-related. This pass should eliminate at least 30% of initial findings. Precision matters more than recall.

When calling store_chunk_findings, provide a JSON array of objects:
[
  {
    "slides": "Slide 5" or "Slides 5, 12",
    "evidence": "exact verbatim text from the slide that contains the error (copy-paste)",
    "issue": "description of the issue — explain WHY the evidence text is wrong",
    "flag": "Critical" or "High" or "Medium" or "Low",
    "remediation": "suggested fix",
    "category": "Spelling" or "Grammar" or "Punctuation" or "Terminology" or "Tone"
  }
]

IMPORTANT: The "evidence" field is REQUIRED. If you cannot provide the exact erroneous text, do not include the finding.

Severity flags:
- Critical: would cause embarrassment or immediate credibility loss
- High: clear errors a customer will notice
- Medium: consistency issues that erode polish
- Low: minor style preferences (use sparingly)

Remediation rules:
- Typos: show corrected word or phrase
- Grammar/Punctuation: show corrected fragment
- Terminology: specify canonical term, suggest find/replace + spot check
- Tone drift: describe the drift in one sentence, suggest target voice (do not rewrite)

VISUAL/LAYOUT REVIEW PROCESS (apply after all text chunks are reviewed):

The analyze_visual_consistency tool returns pre-computed visual/layout findings already formatted for storage.
The result includes a "findings_json" field containing a JSON array of findings ready to store.

Your job is simple:
1. Take the "findings_json" string from the analyze_visual_consistency result.
2. Call store_chunk_findings(chunk_id=-1, findings_json=<the findings_json string>) in a SINGLE call.
   - Do NOT iterate through findings one at a time.
   - Do NOT modify the findings — they are already correctly formatted with slides, evidence, issue, flag, remediation, and category.
   - Pass the entire findings_json string as-is to store_chunk_findings.
3. If the findings_json is "[]" (empty), still call store_chunk_findings(chunk_id=-1, findings_json="[]").

The pre-computed checks cover: font family, font size, text color, bold/italic, shape positioning, shape sizing, fill colors, within-slide font/size/color mixing, clutter, overlap, crammed margins, misaligned columns, and rotation.

For visual/layout findings, use these categories:
- category: "Visual" for consistency issues (fonts, colors, sizes, bold/italic, alignment, fills)
- category: "Layout" for spatial/arrangement issues (clutter, overlap, column alignment, whitespace, rotation)

Visual severity guidelines:
- Critical: Overlapping text making content unreadable
- High: Inconsistent title fonts/sizes/colors across slides (audience will notice); inconsistent bold/italic on titles
- Medium: Body font/size/color deviation; minor alignment drift; inconsistent accent colors; within-slide font/size mixing
- Low: Slightly different margins; minor spacing differences; non-zero rotation on decorative elements

Visual constraints:
- ZERO FALSE POSITIVES applies equally to visual findings.
- Only flag visual issues that are clearly inconsistent — some variation is intentional design.
- Use the dominant pattern (most common font, size, position, color, bold/italic) as the baseline, not an arbitrary standard.
- Do not flag intentional design elements like section dividers or cover slides that are meant to look different.
- In the "evidence" field for visual findings, describe the specific shapes/elements involved and their measured properties (e.g., "Title on Slide 4: Calibri 24pt vs dominant Calibri 32pt on Slides 2, 3, 5-10").

After all reviews are complete, call merge_and_dedupe_findings and present the final results.

CRITICAL TOOL USAGE RULES:
- You MUST call store_chunk_findings for EVERY review chunk. Do NOT skip this tool or summarize findings in text without storing them.
- You MUST call store_chunk_findings with chunk_id=-1 for visual/layout findings. Even if you find no visual issues, store an empty array: store_chunk_findings(chunk_id=-1, findings_json="[]").
- You MUST call merge_and_dedupe_findings after all chunks are stored.
- Do NOT produce your final report until merge_and_dedupe_findings has been called.
"""

# FunctionTool declarations
tool_extract_deck = FunctionTool(
    name="extract_deck",
    description="Extract visible text from a PPTX file, convert to markdown with slide markers.",
    parameters={
        "type": "object",
        "properties": {"pptx_path": {"type": "string", "description": "Path to the PPTX file"}},
        "required": ["pptx_path"], "additionalProperties": False},
    strict=True)

tool_get_review_windows = FunctionTool(
    name="get_review_windows",
    description="Split markdown into overlapping review windows.",
    parameters={
        "type": "object",
        "properties": {
            "markdown": {"type": "string", "description": "Full markdown text"},
            "window_size": {"type": "integer", "description": "Words per window"},
            "overlap": {"type": "integer", "description": "Overlap words"}},
        "required": ["markdown", "window_size", "overlap"], "additionalProperties": False},
    strict=True)

tool_store_chunk_findings = FunctionTool(
    name="store_chunk_findings",
    description="Store QA findings from one chunk.",
    parameters={
        "type": "object",
        "properties": {
            "chunk_id": {"type": "integer", "description": "Chunk ID"},
            "findings_json": {"type": "string", "description": "JSON array of findings"}},
        "required": ["chunk_id", "findings_json"], "additionalProperties": False},
    strict=True)

tool_merge_and_dedupe = FunctionTool(
    name="merge_and_dedupe_findings",
    description="Merge and deduplicate all findings.",
    parameters={"type": "object", "properties": {}, "required": [], "additionalProperties": False},
    strict=True)

tool_extract_deck_visual = FunctionTool(
    name="extract_deck_visual",
    description="Extract raw visual/layout metadata from a PPTX file. Returns per-shape details (fonts, colors, positions, sizes, rotation). Only use when you need raw data; prefer analyze_visual_consistency for the review.",
    parameters={
        "type": "object",
        "properties": {"pptx_path": {"type": "string", "description": "Path to the PPTX file"}},
        "required": ["pptx_path"], "additionalProperties": False},
    strict=True)

tool_analyze_visual = FunctionTool(
    name="analyze_visual_consistency",
    description="Run deterministic visual & layout consistency analysis on a PPTX file. Returns dominant patterns and all deviations: font families, sizes, colors, bold/italic, text alignment, shape positioning, shape sizing, fill colors, within-slide mixing, clutter, overlap, crammed margins, misaligned columns, and rotation. Use this as step 2 of the editorial QA process instead of extract_deck_visual.",
    parameters={
        "type": "object",
        "properties": {"pptx_path": {"type": "string", "description": "Path to the PPTX file"}},
        "required": ["pptx_path"], "additionalProperties": False},
    strict=True)

ALL_TOOLS = [tool_extract_deck, tool_extract_deck_visual, tool_analyze_visual, tool_get_review_windows, tool_store_chunk_findings, tool_merge_and_dedupe]


def list_model_deployments() -> list[str]:
    """Query Foundry project for available model deployments."""
    try:
        credential = AzureCliCredential()
        pc = AIProjectClient(endpoint=ENDPOINT, credential=credential)
        deployments = pc.deployments.list()
        names = []
        for d in deployments:
            if hasattr(d, "name") and d.name:
                if "embedding" in d.name.lower():
                    continue
                names.append(d.name)
        pc.close()
        credential.close()
        if names:
            logger.info(f"Found {len(names)} model deployments: {names}")
            return sorted(names)
    except Exception as e:
        logger.warning(f"Could not list deployments: {e}")
    return [DEFAULT_MODEL, "gpt-4o-mini", "gpt-4.1", "gpt-4.1-mini", "gpt-4.1-nano"]


def init_foundry(model_name: str | None = None, instructions: str | None = None):
    """Initialize the Foundry client and create the agent."""
    model = model_name or DEFAULT_MODEL
    instr = instructions or SYSTEM_INSTRUCTIONS
    credential = AzureCliCredential()
    project_client = AIProjectClient(endpoint=ENDPOINT, credential=credential)
    openai_client = project_client.get_openai_client()
    agent = project_client.agents.create_version(
        agent_name="editorial-qa-reviewer-ui",
        definition=PromptAgentDefinition(
            model=model,
            instructions=instr,
            tools=ALL_TOOLS,
        ),
    )
    logger.info(f"Agent created: {agent.name} v{agent.version} using model '{model}'")
    return project_client, openai_client, agent


# ═══════════════════════════════════════════════════════════════════════════════
# Orchestration
# ═══════════════════════════════════════════════════════════════════════════════

MAX_TOOL_ROUNDS = 50
MAX_RETRIES = 3
RETRY_BASE_DELAY = 2


def execute_tool_call(name: str, arguments: str) -> str:
    func = TOOL_FUNCTIONS.get(name)
    if not func:
        return json.dumps({"error": f"Unknown tool: {name}"})
    try:
        return func(**json.loads(arguments))
    except Exception as e:
        logger.error(f"Tool '{name}' failed: {e}")
        return json.dumps({"error": str(e)})


def call_with_retry(call_fn, max_retries=MAX_RETRIES):
    for attempt in range(max_retries):
        try:
            return call_fn()
        except Exception as e:
            error_str = str(e).lower()
            is_transient = any(kw in error_str for kw in ["rate_limit", "429", "timeout", "503", "502"])
            if is_transient and attempt < max_retries - 1:
                time.sleep(RETRY_BASE_DELAY * (2 ** attempt))
            else:
                raise


def run_editorial_review(pptx_path: str, openai_client, agent, progress_cb=None):
    """Run the full pipeline. progress_cb(message) is called for status updates."""
    findings_accumulator.clear()
    extracted_deck_cache.clear()

    def status(msg):
        logger.info(msg)
        if progress_cb:
            progress_cb(msg)

    user_message = (
        f"Review the PowerPoint deck at path: {pptx_path}\n\n"
        "Follow your instructions exactly — you MUST call each tool:\n"
        "1. Call extract_deck with the PPTX path\n"
        "2. Call analyze_visual_consistency with the same PPTX path\n"
        "3. Call get_review_windows with the markdown\n"
        "4. For EACH chunk: review it, then call store_chunk_findings (MANDATORY — do NOT skip this tool)\n"
        "5. Take the 'findings_json' string from the analyze_visual_consistency result and call store_chunk_findings(chunk_id=-1, findings_json=<that string>) in ONE call (MANDATORY)\n"
        "6. Call merge_and_dedupe_findings\n"
        "7. Present the final editorial QA report\n\n"
        "CRITICAL: You must call store_chunk_findings for EVERY chunk. Do NOT summarize findings in text without storing them first.\n"
        "CRITICAL: Pass the findings_json from analyze_visual_consistency directly to store_chunk_findings — do NOT iterate or modify.\n"
    )

    status(f"Sending initial request to Foundry agent '{agent.name}' v{agent.version}...")
    agent_ref = {"name": agent.name, "version": str(agent.version), "type": "agent_reference"}
    response = call_with_retry(lambda: openai_client.responses.create(
        input=user_message,
        extra_body={"agent_reference": agent_ref},
    ))

    for round_num in range(MAX_TOOL_ROUNDS):
        function_calls = [item for item in response.output if item.type == "function_call"]
        if not function_calls:
            status(f"Agent completed after {round_num} tool rounds")
            break
        status(f"Tool round {round_num + 1}: executing {len(function_calls)} tool call(s)...")
        tool_results = []
        for fc in function_calls:
            status(f"  Running: {fc.name}")
            result = execute_tool_call(fc.name, fc.arguments)
            tool_results.append(FunctionCallOutput(
                type="function_call_output", call_id=fc.call_id, output=result))
        response = call_with_retry(lambda: openai_client.responses.create(
            input=tool_results,
            previous_response_id=response.id,
            extra_body={"agent_reference": agent_ref},
        ))

    final_output = response.output_text if hasattr(response, "output_text") and response.output_text else ""
    if not final_output:
        for item in response.output:
            if hasattr(item, "text"):
                final_output = item.text
                break
    return final_output


# ═══════════════════════════════════════════════════════════════════════════════
# Native folder dialog (Windows)
# ═══════════════════════════════════════════════════════════════════════════════

def open_folder_dialog() -> str:
    """Open native folder browser dialog via tkinter in a subprocess (avoids Gradio thread issues)."""
    try:
        script = (
            "import tkinter as tk;"
            "from tkinter import filedialog;"
            "root = tk.Tk();"
            "root.withdraw();"
            "root.attributes('-topmost', True);"
            "root.focus_force();"
            "root.after(100, lambda: root.focus_force());"
            "path = filedialog.askdirectory(title='Select folder containing PPTX decks', parent=root);"
            "root.destroy();"
            "print(path or '')"
        )
        result = subprocess.run(
            [sys.executable, "-c", script],
            capture_output=True, text=True, timeout=120,
        )
        chosen = result.stdout.strip()
        # tkinter returns forward-slash paths on Windows; normalise
        if chosen:
            chosen = str(pathlib.Path(chosen))
        return chosen
    except Exception as e:
        logger.warning(f"Folder dialog failed: {e}")
        return ""


# ═══════════════════════════════════════════════════════════════════════════════
# UI rendering — high-contrast, theme-safe
# ═══════════════════════════════════════════════════════════════════════════════

SEVERITY_COLORS = {
    "Critical": "#dc2626",
    "High":     "#ea580c",
    "Medium":   "#d97706",
    "Low":      "#2563eb",
}

SEVERITY_BG = {
    "Critical": "#fef2f2",
    "High":     "#fff7ed",
    "Medium":   "#fffbeb",
    "Low":      "#eff6ff",
}

CATEGORY_ICONS = {
    "Spelling":    "\U0001f4dd",
    "Grammar":     "\U0001f4d0",
    "Punctuation": "\u270f\ufe0f",
    "Terminology": "\U0001f516",
    "Tone":        "\U0001f3ad",
    "Visual":      "\U0001f3a8",
    "Layout":      "\U0001f4cf",
}


def _sev_badge(flag: str) -> str:
    c = SEVERITY_COLORS.get(flag, "#6b7280")
    return (
        f'<span style="background:{c};color:#ffffff;padding:3px 12px;'
        f'border-radius:20px;font-size:0.8em;font-weight:700;letter-spacing:0.3px">'
        f'{flag}</span>'
    )


def _html_escape(text: str) -> str:
    """Escape HTML special characters to prevent injection."""
    return (text.replace("&", "&amp;").replace("<", "&lt;")
                .replace(">", "&gt;").replace('"', "&quot;"))


def build_summary_html(all_results: list) -> str:
    """Build a combined summary dashboard across all decks."""
    # Aggregate
    total = 0
    sev_counts = {"Critical": 0, "High": 0, "Medium": 0, "Low": 0}
    cat_counts: Dict[str, int] = {}
    deck_summaries = []

    for dr in all_results:
        findings = dr["findings"]
        total += len(findings)
        deck_sev = {"Critical": 0, "High": 0, "Medium": 0, "Low": 0}
        for f in findings:
            sev = f.get("flag", f.get("severity", "Low"))
            sev_counts[sev] = sev_counts.get(sev, 0) + 1
            deck_sev[sev] = deck_sev.get(sev, 0) + 1
            cat = f.get("category", "Other")
            cat_counts[cat] = cat_counts.get(cat, 0) + 1
        deck_summaries.append({"name": dr["name"], "count": len(findings), "sev": deck_sev})

    # Severity cards
    cards = ""
    for sev in ["Critical", "High", "Medium", "Low"]:
        count = sev_counts[sev]
        color = SEVERITY_COLORS[sev]
        bg = SEVERITY_BG[sev]
        pct = (count / total * 100) if total else 0
        cards += f"""
        <div style="flex:1;min-width:150px;background:{bg};border:2px solid {color};
                     border-radius:12px;padding:20px 16px;text-align:center">
            <div style="font-size:2.2em;font-weight:800;color:{color}">{count}</div>
            <div style="font-size:0.95em;font-weight:600;color:{color};margin-top:2px">{sev}</div>
            <div style="background:#e5e7eb;border-radius:4px;height:6px;margin-top:10px">
                <div style="background:{color};height:6px;border-radius:4px;width:{pct}%"></div>
            </div>
        </div>"""

    # Category breakdown
    cat_rows = ""
    for cat, count in sorted(cat_counts.items(), key=lambda x: -x[1]):
        icon = CATEGORY_ICONS.get(cat, "\U0001f4cb")
        pct = (count / total * 100) if total else 0
        cat_rows += f"""
        <tr style="border-bottom:1px solid #e5e7eb">
            <td style="padding:10px 14px;color:#1f2937;font-weight:500">{icon} {cat}</td>
            <td style="padding:10px 14px;text-align:center;font-weight:700;color:#111827;font-size:1.1em">{count}</td>
            <td style="padding:10px 14px;width:220px">
                <div style="background:#e5e7eb;border-radius:4px;height:8px">
                    <div style="background:#4f46e5;height:8px;border-radius:4px;width:{pct}%"></div>
                </div>
            </td>
        </tr>"""

    # Per-deck summary rows
    deck_rows = ""
    for ds in deck_summaries:
        badges = ""
        for sev in ["Critical", "High", "Medium", "Low"]:
            c = ds["sev"].get(sev, 0)
            if c > 0:
                color = SEVERITY_COLORS[sev]
                badges += (
                    f'<span style="background:{color};color:#fff;padding:2px 8px;'
                    f'border-radius:12px;font-size:0.75em;margin-right:4px">{c} {sev}</span>'
                )
        name_escaped = _html_escape(ds["name"])
        deck_rows += f"""
        <tr style="border-bottom:1px solid #e5e7eb">
            <td style="padding:10px 14px;color:#1f2937;font-weight:600">\U0001f4ca {name_escaped}</td>
            <td style="padding:10px 14px;text-align:center;font-weight:700;color:#111827">{ds['count']}</td>
            <td style="padding:10px 14px">{badges}</td>
        </tr>"""

    deck_table = ""
    if len(all_results) > 1:
        deck_table = f"""
        <div style="margin-top:28px">
            <h3 style="color:#111827;margin-bottom:12px">\U0001f4c1 Per-Deck Breakdown</h3>
            <table style="width:100%;border-collapse:collapse;background:#ffffff;border-radius:8px;overflow:hidden;
                          border:1px solid #e5e7eb">
                <thead><tr style="background:#f1f5f9;border-bottom:2px solid #cbd5e1">
                    <th style="text-align:left;padding:12px 14px;color:#374151;font-weight:700">Deck</th>
                    <th style="text-align:center;padding:12px 14px;color:#374151;font-weight:700">Findings</th>
                    <th style="padding:12px 14px;color:#374151;font-weight:700;text-align:left">Severity</th>
                </tr></thead>
                <tbody>{deck_rows}</tbody>
            </table>
        </div>"""

    # Quick-look findings table (slide + deck + category + severity + issue snippet)
    ql_rows = ""
    for dr in all_results:
        deck_name_ql = _html_escape(dr["name"])
        for f in dr["findings"]:
            slide_str = _html_escape(str(f.get("slides", f.get("slide", "—"))))
            flag = f.get("flag", f.get("severity", "Low"))
            cat = f.get("category", "Other")
            icon = CATEGORY_ICONS.get(cat, "\U0001f4cb")
            issue_text = f.get("issue", "")
            issue_short = _html_escape(issue_text[:90] + ("..." if len(issue_text) > 90 else ""))
            sev_color = SEVERITY_COLORS.get(flag, "#6b7280")
            ql_rows += f"""
            <tr style="border-bottom:1px solid #e5e7eb">
                <td style="padding:8px 12px;color:#1f2937;font-weight:600;white-space:nowrap">{slide_str}</td>
                <td style="padding:8px 12px;color:#374151;max-width:180px;overflow:hidden;text-overflow:ellipsis;white-space:nowrap">{deck_name_ql}</td>
                <td style="padding:8px 12px;color:#374151">{icon} {cat}</td>
                <td style="padding:8px 12px;text-align:center">
                    <span style="background:{sev_color};color:#fff;padding:2px 10px;border-radius:12px;font-size:0.8em;font-weight:600">{flag}</span>
                </td>
                <td style="padding:8px 12px;color:#1f2937;font-size:0.9em">{issue_short}</td>
            </tr>"""

    quick_look = ""
    if ql_rows:
        quick_look = f"""
        <div style="margin-top:28px">
            <h3 style="color:#111827;margin-bottom:12px">\U0001f50d Quick Look — All Findings</h3>
            <div style="overflow-x:auto">
            <table style="width:100%;border-collapse:collapse;background:#ffffff;border-radius:8px;overflow:hidden;
                          border:1px solid #e5e7eb">
                <thead><tr style="background:#f1f5f9;border-bottom:2px solid #cbd5e1">
                    <th style="text-align:left;padding:10px 12px;color:#374151;font-weight:700;white-space:nowrap">Slide</th>
                    <th style="text-align:left;padding:10px 12px;color:#374151;font-weight:700;white-space:nowrap">Deck</th>
                    <th style="text-align:left;padding:10px 12px;color:#374151;font-weight:700">Category</th>
                    <th style="text-align:center;padding:10px 12px;color:#374151;font-weight:700">Severity</th>
                    <th style="text-align:left;padding:10px 12px;color:#374151;font-weight:700">Issue</th>
                </tr></thead>
                <tbody>{ql_rows}</tbody>
            </table>
            </div>
        </div>"""

    return f"""
    <div style="font-family:'Segoe UI',system-ui,-apple-system,sans-serif;color:#111827;
                background:#ffffff;padding:24px;border-radius:16px;border:1px solid #e5e7eb">
        <h2 style="margin:0 0 4px 0;color:#0f172a;font-size:1.5em">\U0001f4ca Editorial QA Summary</h2>
        <p style="color:#64748b;margin:0 0 20px 0;font-size:0.95em">
            Total findings: <strong style="color:#111827;font-size:1.1em">{total}</strong>
            across <strong style="color:#111827">{len(all_results)}</strong> deck(s)
        </p>
        <div style="display:flex;gap:14px;flex-wrap:wrap;margin-bottom:28px">{cards}</div>
        <h3 style="color:#111827;margin-bottom:12px">\U0001f4cb Category Breakdown</h3>
        <table style="width:100%;border-collapse:collapse;background:#ffffff;border-radius:8px;overflow:hidden;
                      border:1px solid #e5e7eb">
            <thead><tr style="background:#f1f5f9;border-bottom:2px solid #cbd5e1">
                <th style="text-align:left;padding:12px 14px;color:#374151;font-weight:700">Category</th>
                <th style="text-align:center;padding:12px 14px;color:#374151;font-weight:700">Count</th>
                <th style="padding:12px 14px;color:#374151;font-weight:700">Distribution</th>
            </tr></thead>
            <tbody>{cat_rows}</tbody>
        </table>
        {quick_look}
        {deck_table}
    </div>"""


def build_findings_html(all_results: list) -> str:
    """Build per-deck, per-slide expandable findings with high-contrast slide headers."""
    all_html = ""

    for dr in all_results:
        findings = dr["findings"]
        slide_blocks = dr["slide_blocks"]
        deck_name = _html_escape(dr["name"])

        # Group findings by slide
        slide_findings: Dict[int, list] = {}
        for f in findings:
            slide_str = str(f.get("slides", f.get("slide", "")))
            nums = re.findall(r"\d+", slide_str)
            for n in nums:
                slide_findings.setdefault(int(n), []).append(f)

        slide_content = {b.slide_num: b.markdown for b in slide_blocks}
        all_slide_nums = sorted(set(list(slide_findings.keys())))

        # Per-slide sections
        slide_sections = ""
        for slide_num in all_slide_nums:
            s_findings = slide_findings.get(slide_num, [])
            content_preview = slide_content.get(slide_num, "")
            content_lines = content_preview.split("\n")
            display_content = _html_escape("\n".join(content_lines[1:]).strip()) if content_lines else ""

            # Finding cards
            finding_cards = ""
            for f in s_findings:
                flag = f.get("flag", f.get("severity", "Low"))
                category = f.get("category", "Other")
                icon = CATEGORY_ICONS.get(category, "\U0001f4cb")
                issue = _html_escape(f.get("issue", "No description"))
                remediation = _html_escape(f.get("remediation", "No remediation provided"))
                evidence = _html_escape(f.get("evidence", ""))
                sev_color = SEVERITY_COLORS.get(flag, "#6b7280")
                sev_bg = SEVERITY_BG.get(flag, "#f9fafb")

                evidence_block = ""
                if evidence:
                    evidence_block = f"""
                    <div style="margin-bottom:10px;background:#fefce8;border-left:4px solid #eab308;
                                padding:10px 14px;border-radius:6px">
                        <div style="font-weight:700;color:#854d0e;margin-bottom:3px;font-size:0.85em;
                                    text-transform:uppercase;letter-spacing:0.5px">Evidence</div>
                        <div style="color:#713f12;line-height:1.5;font-size:0.9em;font-family:'Cascadia Code','Consolas',monospace">{evidence}</div>
                    </div>"""

                finding_cards += f"""
                <div style="background:#ffffff;border:1px solid #d1d5db;border-left:4px solid {sev_color};
                            border-radius:10px;padding:16px;margin-bottom:10px;
                            box-shadow:0 1px 3px rgba(0,0,0,0.06)">
                    <div style="display:flex;justify-content:space-between;align-items:center;margin-bottom:10px">
                        <span style="color:#1f2937;font-weight:700;font-size:0.95em">{icon} {category}</span>
                        {_sev_badge(flag)}
                    </div>
                    <div style="margin-bottom:10px">
                        <div style="font-weight:700;color:#374151;margin-bottom:4px;font-size:0.85em;
                                    text-transform:uppercase;letter-spacing:0.5px">Issue</div>
                        <div style="color:#1f2937;line-height:1.5;font-size:0.95em">{issue}</div>
                    </div>
                    {evidence_block}
                    <div style="background:#ecfdf5;border-left:4px solid #10b981;padding:12px 14px;
                                border-radius:6px">
                        <div style="font-weight:700;color:#065f46;margin-bottom:3px;font-size:0.85em;
                                    text-transform:uppercase;letter-spacing:0.5px">Suggested Fix</div>
                        <div style="color:#047857;line-height:1.5;font-size:0.95em">{remediation}</div>
                    </div>
                </div>"""

            issue_count = len(s_findings)
            # Badge color = worst severity in this slide
            worst = "Low"
            for f in s_findings:
                fl = f.get("flag", "Low")
                if SEVERITY_RANK.get(fl, 0) > SEVERITY_RANK.get(worst, 0):
                    worst = fl
            badge_color = SEVERITY_COLORS.get(worst, "#2563eb")

            slide_sections += f"""
            <details style="margin-bottom:8px;border:1px solid #cbd5e1;border-radius:10px;overflow:hidden;
                            box-shadow:0 1px 3px rgba(0,0,0,0.06)">
                <summary style="padding:14px 20px;
                                background:linear-gradient(135deg, #1e3a5f 0%, #2563eb 100%);
                                color:#ffffff;cursor:pointer;font-weight:700;font-size:1em;
                                display:flex;justify-content:space-between;align-items:center;
                                list-style:none;user-select:none">
                    <span style="color:#ffffff;display:flex;align-items:center;gap:8px">
                        \U0001f4c4 <span style="color:#ffffff;font-size:1.05em">Slide {slide_num}</span>
                    </span>
                    <span style="background:{badge_color};color:#ffffff;padding:3px 14px;border-radius:20px;
                                 font-size:0.8em;font-weight:700">{issue_count} issue{"s" if issue_count != 1 else ""}</span>
                </summary>
                <div style="padding:18px 20px;background:#ffffff">
                    <div style="background:#f8fafc;border:1px solid #e2e8f0;border-radius:8px;padding:14px;
                                margin-bottom:16px;font-family:'Cascadia Code','Consolas',monospace;
                                font-size:0.85em;white-space:pre-wrap;color:#334155;max-height:200px;
                                overflow-y:auto;line-height:1.6">{display_content}</div>
                    <h4 style="margin:0 0 12px 0;color:#1e293b;font-size:0.95em;font-weight:700">
                        Findings ({issue_count})
                    </h4>
                    {finding_cards}
                </div>
            </details>"""

        # Clean slides
        all_extracted = sorted(b.slide_num for b in slide_blocks)
        clean_slides = [n for n in all_extracted if n not in slide_findings]
        clean_note = ""
        if clean_slides:
            clean_list = ", ".join(str(n) for n in clean_slides)
            clean_note = f"""
            <div style="background:#ecfdf5;border:1px solid #a7f3d0;border-radius:10px;
                        padding:14px 18px;margin-top:12px">
                <span style="color:#065f46;font-weight:600">\u2705 Clean slides (no issues):</span>
                <span style="color:#047857"> {clean_list}</span>
            </div>"""

        # Wrap in deck container
        total_issues = len(findings)
        sev_badges = ""
        for sev in ["Critical", "High", "Medium", "Low"]:
            c = sum(1 for f in findings if f.get("flag", f.get("severity", "Low")) == sev)
            if c > 0:
                color = SEVERITY_COLORS[sev]
                sev_badges += (
                    f'<span style="background:{color};color:#fff;padding:3px 10px;'
                    f'border-radius:16px;font-size:0.75em;font-weight:600;margin-left:6px">'
                    f'{c} {sev}</span>'
                )

        # If single deck, skip the deck wrapper header
        if len(all_results) == 1:
            all_html += f"""
            <div style="margin-bottom:20px">
                {slide_sections}
                {clean_note}
            </div>"""
        else:
            all_html += f"""
            <details open style="margin-bottom:20px;border:2px solid #1e40af;border-radius:14px;overflow:hidden">
                <summary style="padding:16px 24px;background:linear-gradient(135deg, #0f172a 0%, #1e3a5f 100%);
                                color:#ffffff;font-weight:800;font-size:1.1em;cursor:pointer;
                                display:flex;justify-content:space-between;align-items:center;
                                list-style:none;user-select:none">
                    <span style="color:#ffffff">\U0001f4ca {deck_name}</span>
                    <span style="display:flex;align-items:center;gap:4px">
                        <span style="color:#94a3b8;font-size:0.8em;margin-right:6px">{total_issues} issues</span>
                        {sev_badges}
                    </span>
                </summary>
                <div style="padding:20px;background:#ffffff">
                    {slide_sections}
                    {clean_note}
                </div>
            </details>"""

    return f"""
    <div style="font-family:'Segoe UI',system-ui,-apple-system,sans-serif;color:#111827;
                background:#ffffff;padding:24px;border-radius:16px;border:1px solid #e5e7eb">
        <h2 style="margin:0 0 6px 0;color:#0f172a;font-size:1.5em">\U0001f50d Per-Slide Findings</h2>
        <p style="color:#64748b;margin:0 0 20px 0;font-size:0.9em">
            Click a slide to expand and see its content + findings
        </p>
        {all_html}
    </div>"""


def build_report_md(all_results: list) -> str:
    """Combine all agent reports into one markdown document."""
    parts = []
    for dr in all_results:
        parts.append(f"## {dr['name']}\n\n{dr['report']}")
    return "\n\n---\n\n".join(parts)


# ═══════════════════════════════════════════════════════════════════════════════
# Foundry state & agent management
# ═══════════════════════════════════════════════════════════════════════════════

_foundry_state = {"project_client": None, "openai_client": None, "agent": None, "model": None, "instructions": None}


def ensure_agent(model_name: str | None = None, instructions: str | None = None):
    """Get or create the Foundry agent. Re-creates if model or instructions change."""
    requested_model = model_name or DEFAULT_MODEL
    requested_instr = (instructions or "").strip() or SYSTEM_INSTRUCTIONS
    needs_recreate = (
        _foundry_state["agent"] is None
        or _foundry_state["model"] != requested_model
        or _foundry_state["instructions"] != requested_instr
    )
    if needs_recreate:
        reason = "first init" if _foundry_state["agent"] is None else (
            f"model changed: '{_foundry_state['model']}' → '{requested_model}'" if _foundry_state["model"] != requested_model
            else "instructions changed"
        )
        logger.info(f"Agent recreation needed: {reason}")
        if _foundry_state["agent"] is not None:
            try:
                old = _foundry_state["agent"]
                _foundry_state["project_client"].agents.delete_version(
                    agent_name=old.name, agent_version=old.version)
                logger.info(f"Deleted old agent version: {old.name} v{old.version}")
            except Exception:
                pass
        pc, oc, ag = init_foundry(requested_model, requested_instr)
        _foundry_state["project_client"] = pc
        _foundry_state["openai_client"] = oc
        _foundry_state["agent"] = ag
        _foundry_state["model"] = requested_model
        _foundry_state["instructions"] = requested_instr
    else:
        logger.info(f"Reusing existing agent: {_foundry_state['agent'].name} v{_foundry_state['agent'].version} (model: {_foundry_state['model']})")
    return _foundry_state["openai_client"], _foundry_state["agent"], _foundry_state["model"]


# ═══════════════════════════════════════════════════════════════════════════════
# Main workflow functions
# ═══════════════════════════════════════════════════════════════════════════════

def _find_pptx_files(folder: str) -> List[pathlib.Path]:
    """Find all .pptx files in a folder (excluding temp files)."""
    p = pathlib.Path(folder)
    files = sorted(p.glob("**/*.pptx"))
    return [f for f in files if not f.name.startswith("~$")]


def _review_single_deck(pptx_path: str, openai_client, agent, progress_cb=None):
    """Review one deck and return its result dict. Returns None if the file can't be processed."""
    try:
        slide_blocks = pptx_to_markdown_slides(pptx_path)
    except Exception as e:
        logger.warning(f"Skipping {pptx_path}: cannot open file — {e}")
        return None
    report = run_editorial_review(pptx_path, openai_client, agent, progress_cb)
    merged = json.loads(merge_and_dedupe_findings())
    return {
        "file": pptx_path,
        "name": pathlib.Path(pptx_path).name,
        "findings": merged["merged_findings"],
        "slide_blocks": slide_blocks,
        "report": report,
        "duplicates_removed": merged.get("duplicates_removed", 0),
    }


def select_folder_and_review(model_name, instructions, progress=gr.Progress(track_tqdm=False)):
    """Open native folder dialog, review ALL decks, return combined results."""
    progress(0, desc="Opening folder selector...")
    folder = open_folder_dialog()
    if not folder:
        return ("", "\u26a0\ufe0f No folder selected — dialog was cancelled.", "", "", "")

    pptx_files = _find_pptx_files(folder)
    if not pptx_files:
        return (
            folder,
            f"\u26a0\ufe0f No `.pptx` files found in: `{_html_escape(folder)}`",
            "", "", "",
        )

    n = len(pptx_files)
    progress(0.02, desc=f"Found {n} deck(s). Initializing agent...")
    openai_client, agent_obj, active_model = ensure_agent(model_name, instructions)

    all_results = []
    skipped = []
    start_all = time.time()
    for i, pptx in enumerate(pptx_files):
        deck_name = pptx.name
        base_pct = 0.05 + (0.85 * i / n)

        def make_cb(dn, bp, total):
            def cb(msg):
                progress(min(bp + 0.8 / total, 0.90), desc=f"[{dn}] {msg}")
            return cb

        progress(base_pct, desc=f"Reviewing {deck_name} ({i+1}/{n})...")
        result = _review_single_deck(str(pptx), openai_client, agent_obj, make_cb(deck_name, base_pct, n))
        if result is None:
            skipped.append(deck_name)
        else:
            all_results.append(result)

    elapsed = time.time() - start_all

    if not all_results:
        skip_list = ", ".join(skipped)
        return (folder, f"\u274c All {n} file(s) failed to open. Skipped: {skip_list}", "", "", "")

    total_findings = sum(len(r["findings"]) for r in all_results)

    progress(0.92, desc="Building reports...")
    summary_html = build_summary_html(all_results)
    findings_html = build_findings_html(all_results)
    report_md = build_report_md(all_results)

    # Build status
    deck_list = ", ".join(f"**{r['name']}** ({len(r['findings'])})" for r in all_results)
    skip_note = ""
    if skipped:
        skip_note = f"\n\n\u26a0\ufe0f Skipped ({len(skipped)}): {', '.join(skipped)}"
    status_md = (
        f"\u2705 **Review complete** in **{elapsed:.1f}s** \u2014 "
        f"**{total_findings}** findings across **{len(all_results)}** deck(s)\n\n"
        f"\U0001f916 Model: `{active_model}` \u2014 "
        f"\U0001f4c1 Folder: `{_html_escape(folder)}`\n\n"
        f"\U0001f4ca Decks: {deck_list}{skip_note}"
    )

    progress(1.0, desc="Done!")
    return (folder, status_md, summary_html, findings_html, report_md)


def review_from_path(folder_path, model_name, instructions, progress=gr.Progress(track_tqdm=False)):
    """Review all decks from a manually-entered folder path."""
    folder = (folder_path or "").strip()
    if not folder:
        return ("", "\u26a0\ufe0f Please enter a folder path.", "", "", "")
    if not os.path.isdir(folder):
        return (folder, f"\u274c Folder not found: `{_html_escape(folder)}`", "", "", "")

    pptx_files = _find_pptx_files(folder)
    if not pptx_files:
        return (
            folder,
            f"\u26a0\ufe0f No `.pptx` files found in: `{_html_escape(folder)}`",
            "", "", "",
        )

    n = len(pptx_files)
    progress(0.02, desc=f"Found {n} deck(s). Initializing agent...")
    openai_client, agent_obj, active_model = ensure_agent(model_name, instructions)

    all_results = []
    skipped = []
    start_all = time.time()
    for i, pptx in enumerate(pptx_files):
        deck_name = pptx.name
        base_pct = 0.05 + (0.85 * i / n)
        progress(base_pct, desc=f"Reviewing {deck_name} ({i+1}/{n})...")
        result = _review_single_deck(str(pptx), openai_client, agent_obj)
        if result is None:
            skipped.append(deck_name)
        else:
            all_results.append(result)

    elapsed = time.time() - start_all

    if not all_results:
        skip_list = ", ".join(skipped)
        return (folder, f"\u274c All {n} file(s) failed to open. Skipped: {skip_list}", "", "", "")

    total_findings = sum(len(r["findings"]) for r in all_results)

    progress(0.92, desc="Building reports...")
    summary_html = build_summary_html(all_results)
    findings_html = build_findings_html(all_results)
    report_md = build_report_md(all_results)

    deck_list = ", ".join(f"**{r['name']}** ({len(r['findings'])})" for r in all_results)
    skip_note = ""
    if skipped:
        skip_note = f"\n\n\u26a0\ufe0f Skipped ({len(skipped)}): {', '.join(skipped)}"
    status_md = (
        f"\u2705 **Review complete** in **{elapsed:.1f}s** \u2014 "
        f"**{total_findings}** findings across **{len(all_results)}** deck(s)\n\n"
        f"\U0001f916 Model: `{active_model}` \u2014 "
        f"\U0001f4c1 Folder: `{_html_escape(folder)}`\n\n"
        f"\U0001f4ca Decks: {deck_list}{skip_note}"
    )

    progress(1.0, desc="Done!")
    return (folder, status_md, summary_html, findings_html, report_md)


def review_uploaded_file(file_obj, model_name, instructions, progress=gr.Progress(track_tqdm=False)):
    """Review a single uploaded PPTX file."""
    if file_obj is None:
        return ("\u26a0\ufe0f Please upload a PPTX file first.", "", "", "")

    pptx_path = file_obj.name if hasattr(file_obj, "name") else str(file_obj)
    if not os.path.exists(pptx_path):
        return (f"\u274c File not found: {pptx_path}", "", "", "")

    progress(0.02, desc="Initializing agent...")
    openai_client, agent_obj, active_model = ensure_agent(model_name, instructions)

    progress(0.05, desc=f"Reviewing deck with model '{active_model}'...")
    start = time.time()
    result = _review_single_deck(pptx_path, openai_client, agent_obj)
    elapsed = time.time() - start

    progress(0.92, desc="Building reports...")
    all_results = [result]
    summary_html = build_summary_html(all_results)
    findings_html = build_findings_html(all_results)
    report_md = build_report_md(all_results)

    status_md = (
        f"\u2705 **Review complete** in **{elapsed:.1f}s** \u2014 "
        f"**{len(result['findings'])}** findings in **{result['name']}**\n\n"
        f"\U0001f916 Model: `{active_model}`"
    )

    progress(1.0, desc="Done!")
    return (status_md, summary_html, findings_html, report_md)


# ═══════════════════════════════════════════════════════════════════════════════
# Gradio Interface
# ═══════════════════════════════════════════════════════════════════════════════

CUSTOM_CSS = """
/* Force full-width layout */
.gradio-container { max-width: 100% !important; width: 100% !important; padding: 0 24px !important; }

/* Make tabs more prominent */
.tab-nav button { font-weight: 600 !important; font-size: 0.95em !important; }
.tab-nav button.selected {
    border-bottom: 3px solid #2563eb !important;
    color: #1e40af !important;
}

/* Details/summary arrow fix for all themes */
details > summary { list-style: none; }
details > summary::-webkit-details-marker { display: none; }
details > summary::before {
    content: "\\25B6  ";
    display: inline-block;
    transition: transform 0.2s ease;
    margin-right: 6px;
    font-size: 0.7em;
    color: #ffffff;
}
details[open] > summary::before { transform: rotate(90deg); }

/* Bigger action buttons */
.action-btn { min-height: 52px !important; font-size: 1.05em !important; }
"""


def refresh_models():
    models = list_model_deployments()
    return gr.update(choices=models, value=models[0] if models else DEFAULT_MODEL)


def build_app():
    available_models = list_model_deployments()

    with gr.Blocks(
        title="Editorial QA Agent \u2014 PPTX Review",
    ) as app:

        # ── Header ──
        gr.HTML("""
        <div style="background:linear-gradient(135deg, #0f172a 0%, #1e3a5f 60%, #2563eb 100%);
                    padding:32px 36px;border-radius:16px;margin-bottom:20px">
            <h1 style="color:#ffffff;margin:0 0 6px 0;font-size:1.8em;font-weight:800">
                \U0001f4cb Editorial QA Agent
            </h1>
            <p style="color:#94a3b8;margin:0;font-size:1em">
                AI-Powered PPTX Editorial Review \u2014 Powered by Microsoft Foundry
            </p>
            <p style="color:#64748b;margin:8px 0 0 0;font-size:0.85em">
                Reviews decks for <strong style="color:#93c5fd">spelling</strong>,
                <strong style="color:#93c5fd">grammar</strong>,
                <strong style="color:#93c5fd">punctuation</strong>,
                <strong style="color:#93c5fd">terminology</strong>,
                <strong style="color:#93c5fd">tone</strong>,
                <strong style="color:#93c5fd">visual consistency</strong>, and
                <strong style="color:#93c5fd">layout quality</strong>
            </p>
        </div>
        """)

        # ── Model selector ──
        with gr.Row():
            model_dropdown = gr.Dropdown(
                label="\U0001f916 Foundry Model",
                choices=available_models,
                value=DEFAULT_MODEL if DEFAULT_MODEL in available_models else (
                    available_models[0] if available_models else ""
                ),
                interactive=True,
                info="Select which model deployment to use for the editorial review",
                scale=6,
            )
            refresh_btn = gr.Button("\U0001f504 Refresh", variant="secondary", size="sm", scale=1)
        refresh_btn.click(fn=refresh_models, inputs=[], outputs=[model_dropdown])

        # ── Agent Instructions (editable) ──
        with gr.Accordion("\u2699\ufe0f Agent Instructions (click to view / edit)", open=False):
            instructions_box = gr.Textbox(
                label="System Instructions",
                value=SYSTEM_INSTRUCTIONS,
                lines=18,
                max_lines=40,
                interactive=True,
                info="Edit the agent's review instructions. Changes take effect on the next review run.",
            )
            with gr.Row():
                reset_instr_btn = gr.Button("\u21a9 Reset to Default", variant="secondary", size="sm")
                instr_status = gr.Markdown(value="", visible=True)
            reset_instr_btn.click(
                fn=lambda: (SYSTEM_INSTRUCTIONS, "*Instructions reset to default.*"),
                inputs=[],
                outputs=[instructions_box, instr_status],
            )

        # ── Hidden state for folder path ──
        folder_state = gr.State("")

        # ── Main action area ──
        gr.HTML("""
        <div style="border-bottom:1px solid #e5e7eb;margin:16px 0"></div>
        """)

        with gr.Row():
            select_btn = gr.Button(
                "\U0001f4c2 Select Folder & Review All Decks",
                variant="primary",
                size="lg",
                elem_classes=["action-btn"],
                scale=2,
            )
            with gr.Column(scale=2):
                folder_input = gr.Textbox(
                    label="Or enter folder path",
                    placeholder=r"C:\Users\you\Documents\Decks",
                )
            path_btn = gr.Button("\u25b6 Review Path", variant="primary", scale=1)
        with gr.Row():
            upload_input = gr.File(label="Or upload a single PPTX", file_types=[".pptx"], scale=3)
            upload_btn = gr.Button("\u25b6 Review Upload", variant="primary", scale=1)

        # ── Status ──
        gr.HTML('<div style="border-bottom:1px solid #e5e7eb;margin:16px 0"></div>')
        status_output = gr.Markdown(label="Status", value="*Select a folder or upload a file to begin.*")

        # ── Results tabs ──
        with gr.Tabs():
            with gr.TabItem("\U0001f4ca Summary Dashboard", id="summary"):
                summary_output = gr.HTML()
            with gr.TabItem("\U0001f50d Per-Slide Findings", id="findings"):
                findings_output = gr.HTML()
            with gr.TabItem("\U0001f4dd Full Agent Report", id="report"):
                report_output = gr.Markdown()

        # ── Wire: Select Folder button ──
        select_btn.click(
            fn=select_folder_and_review,
            inputs=[model_dropdown, instructions_box],
            outputs=[folder_state, status_output, summary_output, findings_output, report_output],
        )

        # ── Wire: Manual path ──
        path_btn.click(
            fn=review_from_path,
            inputs=[folder_input, model_dropdown, instructions_box],
            outputs=[folder_state, status_output, summary_output, findings_output, report_output],
        )

        # ── Wire: Upload ──
        upload_btn.click(
            fn=review_uploaded_file,
            inputs=[upload_input, model_dropdown, instructions_box],
            outputs=[status_output, summary_output, findings_output, report_output],
        )

    return app


# ═══════════════════════════════════════════════════════════════════════════════
# Main
# ═══════════════════════════════════════════════════════════════════════════════

if __name__ == "__main__":
    app = build_app()
    # Use 0.0.0.0 when running in a container (PORT env var set), else localhost
    host = "0.0.0.0" if os.environ.get("PORT") or os.environ.get("CONTAINER") else "127.0.0.1"
    port = int(os.environ.get("PORT", 7860))
    app.launch(
        server_name=host,
        server_port=port,
        share=False,
        inbrowser=(host == "127.0.0.1"),
        css=CUSTOM_CSS,
        theme=gr.themes.Default(),
    )
